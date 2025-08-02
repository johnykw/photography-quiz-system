from flask import Blueprint, request, jsonify, session, send_file
from datetime import datetime
from ..models.quiz import Question, Response, Course, Admin, ScoreSettings, RecommendationSettings, db
import json
from werkzeug.security import check_password_hash
from sqlalchemy import func
import uuid

quiz_bp = Blueprint('quiz', __name__)

@quiz_bp.route('/api/questions', methods=['GET'])
def get_questions():
    questions = Question.query.order_by(Question.order).all()
    return jsonify([{
        'id': q.id,
        'content': q.content,
        'question_type': q.question_type,
        'order': q.order,
        'options': q.options
    } for q in questions])

@quiz_bp.route('/api/submit', methods=['POST'])
def submit_quiz():
    try:
        data = request.json
        if not data or 'answers' not in data:
            return jsonify({'error': '無效的請求數據'}), 400
            
        session_id = str(uuid.uuid4())
        
        total_score = 0
        max_score = 0
        
        for answer in data['answers']:
            if not answer or 'question_id' not in answer or 'answer' not in answer:
                continue
                
            question = Question.query.get(answer['question_id'])
            if not question:
                continue
                
            # 只對技術問題計分（前17題）
            if question.order <= 17:
                max_score += 1
                is_correct = False
                
                if question.question_type == 'single':
                    is_correct = answer['answer'] == question.correct_answer
                elif question.question_type == 'multiple':
                    is_correct = set(answer['answer']) == set(question.correct_answer)
                
                if is_correct:
                    total_score += 1
            else:
                is_correct = None  # 非評分題目
            
            # 保存回應
            response = Response(
                session_id=session_id,
                question_id=question.id,
                answer=answer['answer'],
                is_correct=is_correct
            )
            db.session.add(response)
        
        # 處理"其它"選項的文字輸入
        other_inputs = data.get('other_inputs', {})
        if other_inputs:
            # 將other_inputs保存到數據庫或日誌中
            # 這裡我們可以創建一個新的表來保存這些數據，或者將其添加到現有的Response中
            print(f"用戶自定義輸入: {other_inputs}")
            # 可以考慮將這些數據保存到一個專門的表中，或者作為JSON存儲在Response表的額外字段中
        
        db.session.commit()
        
        # 計算百分比
        percentage = (total_score / max_score * 100) if max_score > 0 else 0
        
        # 使用評分設定計算等級
        user_level = get_user_level_by_score(total_score)
        
        # 設定等級顏色
        level_colors = {
            '攝影新手': '#4CAF50',
            '進階攝影師': '#FF9800', 
            '高階攝影師': '#F44336',
            '中階攝影師': '#FF9800'  # 向後兼容
        }
        
        level = user_level
        level_color = level_colors.get(user_level, '#6c757d')  # 默認灰色
        
        # 獲取推薦課程
        recommended_courses = get_recommended_courses(data['answers'])
        
        return jsonify({
            'session_id': session_id,
            'score': total_score,
            'max_score': max_score,
            'percentage': round(percentage, 1),
            'level': level,
            'level_color': level_color,
            'recommended_courses': recommended_courses
        })
        
    except Exception as e:
        # 記錄詳細錯誤信息
        import traceback
        error_msg = f"提交處理錯誤: {str(e)}\n{traceback.format_exc()}"
        print(error_msg)  # 輸出到控制台
        
        # 回滾數據庫事務
        db.session.rollback()
        
        return jsonify({'error': '服務器內部錯誤，請稍後重試'}), 500

def get_recommended_courses(answers):
    """
    根據分數和興趣智能推薦課程，從數據庫動態讀取課程信息
    只推薦開啟的課程，並基於興趣關聯進行精確配對
    """
    try:
        # 計算分數
        total_score = 0
        max_score = 17  # 前17題為技術問題
        
        for answer in answers:
            if not answer or 'question_id' not in answer:
                continue
                
            question = Question.query.get(answer['question_id'])
            if not question or question.order > 17:
                continue
                
            is_correct = False
            if question.question_type == 'single':
                is_correct = answer['answer'] == question.correct_answer
            elif question.question_type == 'multiple':
                is_correct = set(answer['answer']) == set(question.correct_answer)
            
            if is_correct:
                total_score += 1
        
        # 分析第18題的興趣選擇
        selected_interests = []
        for answer in answers:
            if not answer or 'question_id' not in answer:
                continue
                
            question = Question.query.get(answer['question_id'])
            if not question or question.order != 18:
                continue
                
            # 獲取第18題的選項
            selected_options = answer.get('answer', [])
            if isinstance(selected_options, list):
                for option_index in selected_options:
                    if 0 <= option_index < len(question.options):
                        option_text = question.options[option_index]
                        selected_interests.append(option_text)
        
        # 從數據庫獲取所有開啟的課程
        all_courses = Course.query.filter_by(is_active=True).all()
        
        if not all_courses:
            return get_fallback_courses(total_score, selected_interests)
        
        # 將課程轉換為推薦格式
        def course_to_dict(course, priority=1):
            return {
                'title': course.title,
                'category': course.category,
                'description': course.description,
                'level': course.level,
                'priority': priority
            }
        
        # 開始構建推薦課程列表
        recommended_courses = []
        used_course_ids = set()
        
        # 1. 根據評分設定判斷用戶等級並推薦相應課程
        user_level = get_user_level_by_score(total_score)
        
        # 如果是攝影新手，必須推薦指定的4個課程
        if user_level == '攝影新手':
            beginner_required_titles = [
                'EOS R系列相機全面操作班',
                '基本自動對焦 - 理論班',
                '掌握拍攝設定-拍出準確色彩不求人',
                '鏡頭配搭實用指南'
            ]
            
            priority = 1
            for required_title in beginner_required_titles:
                for course in all_courses:
                    if course.id not in used_course_ids and required_title in course.title:
                        recommended_courses.append(course_to_dict(course, priority))
                        used_course_ids.add(course.id)
                        priority += 1
                        break
        
        # 2. 基於興趣關聯的智能配對
        for interest in selected_interests:
            for course in all_courses:
                if course.id in used_course_ids:
                    continue
                
                # 檢查課程的興趣標籤是否匹配
                course_tags = []
                if hasattr(course, 'interest_tags') and course.interest_tags:
                    if isinstance(course.interest_tags, list):
                        course_tags = course.interest_tags
                    else:
                        try:
                            course_tags = json.loads(course.interest_tags)
                        except:
                            course_tags = []
                
                # 如果課程的興趣標籤包含用戶選擇的興趣
                if interest in course_tags:
                    priority = len(recommended_courses) + 1
                    recommended_courses.append(course_to_dict(course, priority))
                    used_course_ids.add(course.id)
                    
                    # 限制每個興趣最多推薦4個課程
                    interest_count = sum(1 for c in recommended_courses if interest in course_tags)
                    if interest_count >= 4:
                        break
        
        # 3. 如果推薦課程不足，補充其他開啟的課程
        # 獲取推薦設定
        recommendation_setting = get_active_recommendation_setting()
        min_courses = recommendation_setting['min_courses']
        max_courses = recommendation_setting['max_courses']
        
        if len(recommended_courses) < min_courses:
            remaining_courses = [c for c in all_courses if c.id not in used_course_ids]
            needed = min(max_courses - len(recommended_courses), len(remaining_courses))
            
            for course in remaining_courses[:needed]:
                priority = len(recommended_courses) + 1
                recommended_courses.append(course_to_dict(course, priority))
                used_course_ids.add(course.id)
        
        # 4. 限制推薦數量根據設定
        recommended_courses = recommended_courses[:max_courses]
        
        # 5. 按優先級排序
        recommended_courses.sort(key=lambda x: x['priority'])
        
        return recommended_courses
        
    except Exception as e:
        print(f"課程推薦錯誤: {e}")
        return get_fallback_courses(total_score, selected_interests)

def get_fallback_courses(total_score, selected_interests):
    """
    當數據庫查詢失敗時的備用課程推薦
    """
    # 根據評分設定判斷用戶等級
    user_level = get_user_level_by_score(total_score)
    
    # 新手必推課程
    beginner_courses = [
        {
            'title': '【新手入門】EOS R系列相機全面操作班',
            'category': '新手入門',
            'description': '全面學習EOS R系列相機的操作技巧',
            'level': '攝影新手',
            'priority': 1
        },
        {
            'title': '【新手入門】基本自動對焦 - 理論班',
            'category': '新手入門', 
            'description': '掌握自動對焦的基本理論和應用',
            'level': '攝影新手',
            'priority': 2
        },
        {
            'title': '【新手入門】掌握拍攝設定-拍出準確色彩不求人',
            'category': '新手入門',
            'description': '學習正確的拍攝設定，拍出準確色彩',
            'level': '攝影新手',
            'priority': 3
        },
        {
            'title': '【新手入門】鏡頭配搭實用指南',
            'category': '新手入門',
            'description': '了解不同鏡頭的特性和配搭技巧',
            'level': '攝影新手',
            'priority': 4
        }
    ]
    
    # 人像攝影課程
    portrait_courses = [
        {
            'title': '【新手入門】日常人像生活拍攝攻略',
            'category': '新手入門',
            'description': '學習日常人像攝影技巧',
            'level': '攝影新手',
            'priority': 5
        },
        {
            'title': '【進階攝影】人像攝影技能全面解鎖工作坊 (Cherry Wong)',
            'category': '進階攝影',
            'description': '全面提升人像攝影技能',
            'level': '進階攝影師',
            'priority': 6
        }
    ]
    
    # 舞台攝影課程
    stage_courses = [
        {
            'title': '【新手入門】追星拍攝攻略',
            'category': '新手入門',
            'description': '學習演唱會和舞台攝影技巧',
            'level': '攝影新手',
            'priority': 5
        },
        {
            'title': '【進階攝影】攝動定格：應援攝影工作坊(Cherry Wong)',
            'category': '進階攝影',
            'description': '專業應援攝影技巧',
            'level': '進階攝影師',
            'priority': 6
        }
    ]
    
    recommended_courses = []
    
    # 根據評分設定推薦課程
    if user_level == '攝影新手':
        recommended_courses.extend(beginner_courses)
    
    # 根據興趣添加相關課程
    has_portrait_interest = any('人物攝影' in interest or '人像攝影' in interest for interest in selected_interests)
    has_stage_interest = any('舞台攝影' in interest for interest in selected_interests)
    
    if has_portrait_interest:
        recommended_courses.extend(portrait_courses)
    
    if has_stage_interest:
        recommended_courses.extend(stage_courses)
    
    # 獲取推薦設定
    recommendation_setting = get_active_recommendation_setting()
    min_courses = recommendation_setting['min_courses']
    max_courses = recommendation_setting['max_courses']
    
    # 確保至少有設定的最少課程數量
    if len(recommended_courses) < min_courses:
        recommended_courses.extend(beginner_courses[:min_courses-len(recommended_courses)])
    
    # 限制最多課程數量根據設定
    recommended_courses = recommended_courses[:max_courses]
    
    return recommended_courses

# 統計分析API端點

@quiz_bp.route('/api/admin/login', methods=['POST'])
def admin_login():
    data = request.json
    admin = Admin.query.filter_by(username=data['username']).first()
    
    if admin and check_password_hash(admin.password_hash, data['password']):
        session['admin_logged_in'] = True
        session['admin_id'] = admin.id
        return jsonify({'success': True})
    
    return jsonify({'success': False, 'message': '用戶名或密碼錯誤'}), 401

@quiz_bp.route('/api/admin/logout', methods=['POST'])
def admin_logout():
    session.pop('admin_logged_in', None)
    session.pop('admin_id', None)
    return jsonify({'success': True})

@quiz_bp.route('/api/admin/stats', methods=['GET'])
def get_admin_stats():
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    total_responses = db.session.query(Response.session_id).distinct().count()
    total_questions = Question.query.count()
    
    # 計算平均分數
    scores = db.session.query(
        Response.session_id,
        func.sum(Response.is_correct.cast(db.Integer)).label('score')
    ).filter(
        Response.is_correct.isnot(None)
    ).group_by(Response.session_id).all()
    
    avg_score = sum(score.score for score in scores) / len(scores) if scores else 0
    
    return jsonify({
        'total_responses': total_responses,
        'total_questions': total_questions,
        'avg_score': round(avg_score, 1)
    })

@quiz_bp.route('/api/admin/real_time_stats', methods=['GET'])
def get_real_time_stats():
    """獲取即時統計數據（不顯示正確答案）"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    # 獲取所有回應數據
    responses = Response.query.all()
    
    # 總回應數
    total_responses = len(set(r.session_id for r in responses))
    
    # 問題統計
    questions = Question.query.order_by(Question.order).all()
    question_stats = []
    
    for question in questions:
        question_responses = [r for r in responses if r.question_id == question.id]
        total_answers = len(question_responses)
        
        if question.order <= 17:  # 技術問題
            correct_answers = len([r for r in question_responses if r.is_correct])
            correct_rate = (correct_answers / total_answers * 100) if total_answers > 0 else 0
        else:
            correct_answers = 0
            correct_rate = 0
        
        # 選項統計
        option_stats = []
        if total_answers > 0:
            for i, option in enumerate(question.options):
                if question.question_type == 'single':
                    count = len([r for r in question_responses if r.answer == i])
                else:  # multiple choice
                    count = len([r for r in question_responses if i in (r.answer or [])])
                
                percentage = (count / total_answers * 100) if total_answers > 0 else 0
                option_stats.append({
                    'option': option,
                    'count': count,
                    'percentage': round(percentage, 1)
                })
        
        question_stats.append({
            'id': question.id,
            'order': question.order,
            'content': question.content,
            'question_type': question.question_type,
            'correct_rate': round(correct_rate, 1),
            'correct_answers': correct_answers,
            'total_answers': total_answers,
            'option_stats': option_stats
        })
    
    return jsonify({
        'total_responses': total_responses,
        'question_stats': question_stats
    })


@quiz_bp.route('/api/admin/detailed_stats', methods=['GET'])
def get_detailed_stats():
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    start_date = request.args.get('start_date')
    end_date = request.args.get('end_date')
    
    # 基礎查詢
    query = Response.query
    
    # 日期篩選
    if start_date:
        query = query.filter(Response.created_at >= datetime.fromisoformat(start_date))
    if end_date:
        query = query.filter(Response.created_at <= datetime.fromisoformat(end_date))
    
    responses = query.all()
    
    # 總回應數
    total_responses = len(set(r.session_id for r in responses))
    
    # 問題統計
    questions = Question.query.order_by(Question.order).all()
    question_stats = []
    
    for question in questions:
        question_responses = [r for r in responses if r.question_id == question.id]
        total_answers = len(question_responses)
        
        if question.order <= 17:  # 技術問題
            correct_answers = len([r for r in question_responses if r.is_correct])
            correct_rate = (correct_answers / total_answers * 100) if total_answers > 0 else 0
        else:
            correct_answers = 0
            correct_rate = 0
        
        # 選項統計
        option_stats = []
        if total_answers > 0:
            for i, option in enumerate(question.options):
                if question.question_type == 'single':
                    count = len([r for r in question_responses if r.answer == i])
                else:  # multiple choice
                    count = len([r for r in question_responses if i in (r.answer or [])])
                
                percentage = (count / total_answers * 100) if total_answers > 0 else 0
                option_stats.append({
                    'option': option,
                    'count': count,
                    'percentage': round(percentage, 1)
                })
        
        question_stats.append({
            'id': question.id,
            'order': question.order,
            'content': question.content,
            'question_type': question.question_type,
            'correct_rate': round(correct_rate, 1),
            'correct_answers': correct_answers,
            'total_answers': total_answers,
            'option_stats': option_stats
        })
    
    # 分數分布統計
    score_distribution = []
    session_scores = db.session.query(
        Response.session_id,
        func.sum(Response.is_correct.cast(db.Integer)).label('score')
    ).filter(
        Response.is_correct.isnot(None)
    ).group_by(Response.session_id)
    
    if start_date or end_date:
        if start_date:
            session_scores = session_scores.filter(Response.created_at >= datetime.fromisoformat(start_date))
        if end_date:
            session_scores = session_scores.filter(Response.created_at <= datetime.fromisoformat(end_date))
    
    scores = [s.score for s in session_scores.all()]
    
    for score in range(18):  # 0-17分
        count = scores.count(score)
        percentage = (count / len(scores) * 100) if scores else 0
        score_distribution.append({
            'score': score,
            'count': count,
            'percentage': round(percentage, 1)
        })
    
    return jsonify({
        'total_responses': total_responses,
        'question_stats': question_stats,
        'score_distribution': score_distribution
    })

@quiz_bp.route('/api/admin/clear_data', methods=['POST'])
def clear_data():
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    data = request.json
    start_date = data.get('start_date')
    end_date = data.get('end_date')
    clear_all = data.get('clear_all', False)
    
    if clear_all:
        Response.query.delete()
    else:
        query = Response.query
        if start_date:
            query = query.filter(Response.created_at >= datetime.fromisoformat(start_date))
        if end_date:
            query = query.filter(Response.created_at <= datetime.fromisoformat(end_date))
        query.delete()
    
    db.session.commit()
    return jsonify({'success': True})


# 問題管理API端點

@quiz_bp.route('/api/admin/questions', methods=['GET'])
def get_admin_questions():
    """獲取所有問題（管理員用）"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    questions = Question.query.order_by(Question.order).all()
    return jsonify([{
        'id': q.id,
        'content': q.content,
        'question_type': q.question_type,
        'order': q.order,
        'options': q.options,
        'correct_answer': q.correct_answer,
        'created_at': q.created_at.isoformat() if q.created_at else None
    } for q in questions])

@quiz_bp.route('/api/admin/questions', methods=['POST'])
def add_question():
    """添加新問題"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    data = request.json
    
    # 獲取下一個順序號
    max_order = db.session.query(func.max(Question.order)).scalar() or 0
    
    question = Question(
        content=data['content'],
        question_type=data['question_type'],
        order=max_order + 1,
        options=data['options'],
        correct_answer=data.get('correct_answer')
    )
    
    db.session.add(question)
    db.session.commit()
    
    return jsonify({
        'success': True,
        'question': {
            'id': question.id,
            'content': question.content,
            'question_type': question.question_type,
            'order': question.order,
            'options': question.options,
            'correct_answer': question.correct_answer
        }
    })

@quiz_bp.route('/api/admin/questions/<int:question_id>', methods=['PUT'])
def update_question(question_id):
    """更新問題"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    question = Question.query.get_or_404(question_id)
    data = request.json
    
    question.content = data['content']
    question.question_type = data['question_type']
    question.options = data['options']
    question.correct_answer = data.get('correct_answer')
    
    db.session.commit()
    
    return jsonify({
        'success': True,
        'question': {
            'id': question.id,
            'content': question.content,
            'question_type': question.question_type,
            'order': question.order,
            'options': question.options,
            'correct_answer': question.correct_answer
        }
    })

@quiz_bp.route('/api/admin/questions/<int:question_id>', methods=['DELETE'])
def delete_question(question_id):
    """刪除問題"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    question = Question.query.get_or_404(question_id)
    
    # 刪除相關的回應記錄
    Response.query.filter_by(question_id=question_id).delete()
    
    # 刪除問題
    db.session.delete(question)
    db.session.commit()
    
    return jsonify({'success': True})

@quiz_bp.route('/api/admin/questions/reorder', methods=['POST'])
def reorder_questions():
    """重新排序問題"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    data = request.json
    question_orders = data['questions']  # [{'id': 1, 'order': 1}, ...]
    
    for item in question_orders:
        question = Question.query.get(item['id'])
        if question:
            question.order = item['order']
    
    db.session.commit()
    return jsonify({'success': True})

# 課程管理API端點

@quiz_bp.route('/api/admin/courses', methods=['GET'])
def get_admin_courses():
    """獲取所有課程（管理員用）"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    courses = Course.query.all()
    return jsonify([{
        'id': c.id,
        'title': c.title,
        'description': c.description,
        'category': c.category,
        'level': c.level,
        'is_active': getattr(c, 'is_active', True),
        'interest_tags': getattr(c, 'interest_tags', []) or [],
        'created_at': c.created_at.isoformat() if c.created_at else None
    } for c in courses])

@quiz_bp.route('/api/admin/courses', methods=['POST'])
def add_course():
    """添加新課程"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    data = request.json
    
    course = Course(
        title=data['title'],
        description=data['description'],
        category=data['category'],
        level=data['level']
    )
    
    # 設置新欄位
    if hasattr(course, 'is_active'):
        course.is_active = data.get('is_active', True)
    if hasattr(course, 'interest_tags'):
        course.interest_tags = json.dumps(data.get('interest_tags', []))
    
    db.session.add(course)
    db.session.commit()
    
    return jsonify({
        'success': True,
        'course': {
            'id': course.id,
            'title': course.title,
            'description': course.description,
            'category': course.category,
            'level': course.level,
            'is_active': getattr(course, 'is_active', True),
            'interest_tags': getattr(course, 'interest_tags', []) or []
        }
    })

@quiz_bp.route('/api/admin/courses/<int:course_id>', methods=['PUT'])
def update_course(course_id):
    """更新課程"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    course = Course.query.get_or_404(course_id)
    data = request.json
    
    course.title = data['title']
    course.description = data['description']
    course.category = data['category']
    course.level = data['level']
    
    # 更新新欄位
    if hasattr(course, 'is_active'):
        course.is_active = data.get('is_active', True)
    if hasattr(course, 'interest_tags'):
        course.interest_tags = json.dumps(data.get('interest_tags', []))
    
    db.session.commit()
    
    return jsonify({
        'success': True,
        'course': {
            'id': course.id,
            'title': course.title,
            'description': course.description,
            'category': course.category,
            'level': course.level,
            'is_active': getattr(course, 'is_active', True),
            'interest_tags': getattr(course, 'interest_tags', []) or []
        }
    })

@quiz_bp.route('/api/admin/courses/<int:course_id>', methods=['DELETE'])
def delete_course(course_id):
    """刪除課程"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    course = Course.query.get_or_404(course_id)
    db.session.delete(course)
    db.session.commit()
    
    return jsonify({'success': True})

# 用戶資料管理API端點

@quiz_bp.route('/api/admin/profile', methods=['GET'])
def get_admin_profile():
    """獲取管理員資料"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    admin = Admin.query.get(session['admin_id'])
    if not admin:
        return jsonify({'error': '管理員不存在'}), 404
    
    return jsonify({
        'id': admin.id,
        'username': admin.username,
        'created_at': admin.created_at.isoformat() if admin.created_at else None
    })

@quiz_bp.route('/api/admin/profile', methods=['PUT'])
def update_admin_profile():
    """更新管理員資料"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    admin = Admin.query.get(session['admin_id'])
    if not admin:
        return jsonify({'error': '管理員不存在'}), 404
    
    data = request.json
    
    # 更新用戶名
    if 'username' in data:
        # 檢查用戶名是否已存在
        existing_admin = Admin.query.filter_by(username=data['username']).first()
        if existing_admin and existing_admin.id != admin.id:
            return jsonify({'error': '用戶名已存在'}), 400
        admin.username = data['username']
    
    # 更新密碼
    if 'password' in data and data['password']:
        from werkzeug.security import generate_password_hash
        admin.password_hash = generate_password_hash(data['password'])
    
    db.session.commit()
    
    return jsonify({
        'success': True,
        'admin': {
            'id': admin.id,
            'username': admin.username
        }
    })



# 數據導出API端點

@quiz_bp.route('/api/admin/export/excel', methods=['GET', 'POST'])
def export_excel():
    """導出Excel格式的統計數據"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.chart import BarChart, Reference
        import io
        import base64
        
        # 處理GET和POST請求
        if request.method == 'POST':
            data = request.json or {}
        else:
            data = request.args
        
        start_date = data.get('start_date')
        end_date = data.get('end_date')
        
        # 獲取篩選後的數據
        query = Response.query
        if start_date:
            query = query.filter(Response.created_at >= datetime.fromisoformat(start_date))
        if end_date:
            query = query.filter(Response.created_at <= datetime.fromisoformat(end_date))
        
        responses = query.all()
        
        # 創建Excel工作簿
        wb = openpyxl.Workbook()
        
        # 總覽工作表
        ws_summary = wb.active
        ws_summary.title = "統計總覽"
        
        # 設置標題樣式
        title_font = Font(size=16, bold=True)
        header_font = Font(size=12, bold=True)
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        # 總覽數據
        total_responses = len(set(r.session_id for r in responses))
        questions = Question.query.order_by(Question.order).all()
        total_questions = len(questions)
        
        # 計算平均分數
        session_scores = {}
        for response in responses:
            if response.session_id not in session_scores:
                session_scores[response.session_id] = {'correct': 0, 'total': 0}
            if response.question_id <= 17:  # 技術問題
                session_scores[response.session_id]['total'] += 1
                if response.is_correct:
                    session_scores[response.session_id]['correct'] += 1
        
        avg_score = sum(s['correct'] for s in session_scores.values()) / len(session_scores) if session_scores else 0
        
        # 寫入總覽數據
        ws_summary['A1'] = "攝影問卷系統統計報告"
        ws_summary['A1'].font = title_font
        ws_summary['A3'] = "統計期間："
        ws_summary['B3'] = f"{start_date or '開始'} 至 {end_date or '現在'}"
        ws_summary['A4'] = "總回應數："
        ws_summary['B4'] = total_responses
        ws_summary['A5'] = "問題總數："
        ws_summary['B5'] = total_questions
        ws_summary['A6'] = "平均分數："
        ws_summary['B6'] = f"{avg_score:.1f}"
        
        # 新增：參與者分數統計
        ws_summary['A8'] = "參與者分數分布統計"
        ws_summary['A8'].font = header_font
        
        # 計算分數分布
        score_distribution = {}
        for session_id, score_data in session_scores.items():
            score = score_data['correct']
            if score not in score_distribution:
                score_distribution[score] = 0
            score_distribution[score] += 1
        
        # 寫入分數分布統計
        row = 9
        ws_summary[f'A{row}'] = "分數"
        ws_summary[f'B{row}'] = "人數"
        ws_summary[f'C{row}'] = "百分比"
        
        # 設置表頭樣式
        for col in ['A', 'B', 'C']:
            cell = ws_summary[f'{col}{row}']
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
        
        row += 1
        for score in sorted(score_distribution.keys()):
            count = score_distribution[score]
            percentage = (count / total_responses * 100) if total_responses > 0 else 0
            
            ws_summary[f'A{row}'] = f"{score}分"
            ws_summary[f'B{row}'] = count
            ws_summary[f'C{row}'] = f"{percentage:.1f}%"
            
            # 設置對齊方式
            ws_summary[f'A{row}'].alignment = Alignment(horizontal='center')
            ws_summary[f'B{row}'].alignment = Alignment(horizontal='center')
            ws_summary[f'C{row}'].alignment = Alignment(horizontal='center')
            
            row += 1
        
        # 添加總計行
        ws_summary[f'A{row}'] = "總計"
        ws_summary[f'A{row}'].font = header_font
        ws_summary[f'B{row}'] = total_responses
        ws_summary[f'B{row}'].font = header_font
        ws_summary[f'C{row}'] = "100.0%"
        ws_summary[f'C{row}'].font = header_font
        
        # 設置總計行樣式
        for col in ['A', 'B', 'C']:
            cell = ws_summary[f'{col}{row}']
            cell.alignment = Alignment(horizontal='center')
            cell.fill = PatternFill(start_color="E7E6E6", end_color="E7E6E6", fill_type="solid")
        
        # 詳細統計工作表
        ws_detail = wb.create_sheet("詳細統計")
        
        # 設置表頭
        headers = ['問題編號', '問題內容', '問題類型', '總回答數', '正確答案數', '正確率(%)', '選項1', '選項1人數', '選項1比例(%)', '選項2', '選項2人數', '選項2比例(%)', '選項3', '選項3人數', '選項3比例(%)', '選項4', '選項4人數', '選項4比例(%)']
        
        for col, header in enumerate(headers, 1):
            cell = ws_detail.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
        
        # 填入問題統計數據
        for row, question in enumerate(questions, 2):
            question_responses = [r for r in responses if r.question_id == question.id]
            total_answers = len(question_responses)
            
            if question.order <= 17:  # 技術問題
                correct_answers = len([r for r in question_responses if r.is_correct])
                correct_rate = (correct_answers / total_answers * 100) if total_answers > 0 else 0
            else:
                correct_answers = 0
                correct_rate = 0
            
            ws_detail.cell(row=row, column=1, value=question.order)
            ws_detail.cell(row=row, column=2, value=question.content)
            ws_detail.cell(row=row, column=3, value='單選題' if question.question_type == 'single' else '多選題')
            ws_detail.cell(row=row, column=4, value=total_answers)
            ws_detail.cell(row=row, column=5, value=correct_answers)
            ws_detail.cell(row=row, column=6, value=f"{correct_rate:.1f}")
            
            # 選項統計
            for i, option in enumerate(question.options[:4]):  # 最多4個選項
                if question.question_type == 'single':
                    count = len([r for r in question_responses if r.answer == i])
                else:
                    count = len([r for r in question_responses if i in (r.answer or [])])
                
                percentage = (count / total_answers * 100) if total_answers > 0 else 0
                
                ws_detail.cell(row=row, column=7 + i*3, value=option)
                ws_detail.cell(row=row, column=8 + i*3, value=count)
                ws_detail.cell(row=row, column=9 + i*3, value=f"{percentage:.1f}")
        
        # 調整列寬
        for column in ws_detail.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws_detail.column_dimensions[column_letter].width = adjusted_width
        
        # 保存到內存
        excel_buffer = io.BytesIO()
        wb.save(excel_buffer)
        excel_buffer.seek(0)
        
        # 轉換為base64
        excel_data = base64.b64encode(excel_buffer.getvalue()).decode()
        
        return jsonify({
            'success': True,
            'data': excel_data,
            'filename': f'攝影問卷統計_{datetime.now().strftime("%Y%m%d_%H%M%S")}.xlsx'
        })
        
    except Exception as e:
        return jsonify({'error': f'導出Excel失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/export/powerpoint', methods=['GET', 'POST'])
def export_powerpoint():
    """導出PowerPoint格式的統計數據"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import matplotlib.pyplot as plt
        import seaborn as sns
        import io
        import base64
        import tempfile
        import os
        
        # 處理GET和POST請求
        if request.method == 'POST':
            data = request.json or {}
        else:
            data = request.args
            
        start_date = data.get('start_date')
        end_date = data.get('end_date')
        
        # 獲取篩選後的數據
        query = Response.query
        if start_date:
            query = query.filter(Response.created_at >= datetime.fromisoformat(start_date))
        if end_date:
            query = query.filter(Response.created_at <= datetime.fromisoformat(end_date))
        
        responses = query.all()
        questions = Question.query.order_by(Question.order).all()
        
        # 創建PowerPoint演示文稿
        prs = Presentation()
        
        # 設置中文字體
        plt.rcParams['font.sans-serif'] = ['DejaVu Sans', 'SimHei', 'Arial Unicode MS']
        plt.rcParams['axes.unicode_minus'] = False
        
        # 第一張幻燈片：標題頁
        slide_layout = prs.slide_layouts[0]  # 標題幻燈片
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "攝影問卷系統統計報告"
        subtitle.text = f"統計期間：{start_date or '開始'} 至 {end_date or '現在'}"
        
        # 第二張幻燈片：總覽統計
        slide_layout = prs.slide_layouts[1]  # 標題和內容
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "統計總覽"
        
        # 計算統計數據
        total_responses = len(set(r.session_id for r in responses))
        total_questions = len(questions)
        
        session_scores = {}
        for response in responses:
            if response.session_id not in session_scores:
                session_scores[response.session_id] = {'correct': 0, 'total': 0}
            if response.question_id <= 17:  # 技術問題
                session_scores[response.session_id]['total'] += 1
                if response.is_correct:
                    session_scores[response.session_id]['correct'] += 1
        
        avg_score = sum(s['correct'] for s in session_scores.values()) / len(session_scores) if session_scores else 0
        
        # 添加文本框
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = f"總回應數：{total_responses}"
        p.font.size = Pt(24)
        
        p = text_frame.add_paragraph()
        p.text = f"問題總數：{total_questions}"
        p.font.size = Pt(24)
        
        p = text_frame.add_paragraph()
        p.text = f"平均分數：{avg_score:.1f}"
        p.font.size = Pt(24)
        
        # 第三張幻燈片：正確率圖表
        slide_layout = prs.slide_layouts[5]  # 空白幻燈片
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加標題
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "各題正確率統計"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # 創建正確率圖表
        question_numbers = []
        correct_rates = []
        
        for question in questions[:17]:  # 只統計技術問題
            question_responses = [r for r in responses if r.question_id == question.id]
            total_answers = len(question_responses)
            correct_answers = len([r for r in question_responses if r.is_correct])
            correct_rate = (correct_answers / total_answers * 100) if total_answers > 0 else 0
            
            question_numbers.append(f"Q{question.order}")
            correct_rates.append(correct_rate)
        
        # 創建圖表
        plt.figure(figsize=(12, 6))
        bars = plt.bar(question_numbers, correct_rates, color='#4472C4', alpha=0.8)
        plt.title('各題正確率統計', fontsize=16, fontweight='bold')
        plt.xlabel('問題編號', fontsize=12)
        plt.ylabel('正確率 (%)', fontsize=12)
        plt.ylim(0, 100)
        
        # 在柱狀圖上添加數值標籤
        for bar, rate in zip(bars, correct_rates):
            plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1, 
                    f'{rate:.1f}%', ha='center', va='bottom', fontsize=10)
        
        plt.xticks(rotation=45)
        plt.tight_layout()
        
        # 保存圖表到臨時文件
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            plt.savefig(tmp_file.name, dpi=300, bbox_inches='tight')
            chart_path = tmp_file.name
        plt.close()
        
        # 添加圖表到幻燈片
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        slide.shapes.add_picture(chart_path, left, top, width, height)
        
        # 清理臨時文件
        os.unlink(chart_path)
        
        # 第四張幻燈片：回應分布圖
        slide_layout = prs.slide_layouts[5]  # 空白幻燈片
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加標題
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "分數分布統計"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # 創建分數分布圖
        scores = [s['correct'] for s in session_scores.values()]
        
        plt.figure(figsize=(10, 6))
        plt.hist(scores, bins=range(0, max(scores)+2), color='#70AD47', alpha=0.8, edgecolor='black')
        plt.title('分數分布統計', fontsize=16, fontweight='bold')
        plt.xlabel('正確答題數', fontsize=12)
        plt.ylabel('人數', fontsize=12)
        plt.grid(axis='y', alpha=0.3)
        
        # 保存圖表到臨時文件
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            plt.savefig(tmp_file.name, dpi=300, bbox_inches='tight')
            chart_path = tmp_file.name
        plt.close()
        
        # 添加圖表到幻燈片
        slide.shapes.add_picture(chart_path, left, top, width, height)
        
        # 清理臨時文件
        os.unlink(chart_path)
        
        # 保存PowerPoint到內存
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        # 轉換為base64
        ppt_data = base64.b64encode(ppt_buffer.getvalue()).decode()
        
        return jsonify({
            'success': True,
            'data': ppt_data,
            'filename': f'攝影問卷統計_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
        })
        
    except Exception as e:
        return jsonify({'error': f'導出PowerPoint失敗: {str(e)}'}), 500



# ==================== 評分設定管理API ====================

@quiz_bp.route('/api/admin/score-settings', methods=['GET'])
def get_score_settings():
    """獲取所有評分設定"""
    try:
        settings = ScoreSettings.query.order_by(ScoreSettings.order).all()
        
        settings_data = []
        for setting in settings:
            settings_data.append({
                'id': setting.id,
                'level_name': setting.level_name,
                'min_score': setting.min_score,
                'max_score': setting.max_score,
                'description': setting.description,
                'is_active': setting.is_active,
                'order': setting.order,
                'created_at': setting.created_at.isoformat() if setting.created_at else None,
                'updated_at': setting.updated_at.isoformat() if setting.updated_at else None
            })
        
        return jsonify({
            'success': True,
            'settings': settings_data
        })
        
    except Exception as e:
        return jsonify({'error': f'獲取評分設定失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/score-settings/<int:setting_id>', methods=['PUT'])
def update_score_setting(setting_id):
    """更新評分設定"""
    try:
        data = request.json
        if not data:
            return jsonify({'error': '無效的請求數據'}), 400
        
        setting = ScoreSettings.query.get(setting_id)
        if not setting:
            return jsonify({'error': '評分設定不存在'}), 404
        
        # 驗證分數範圍
        min_score = data.get('min_score', setting.min_score)
        max_score = data.get('max_score', setting.max_score)
        
        if min_score >= max_score:
            return jsonify({'error': '最低分數必須小於最高分數'}), 400
        
        # 檢查分數範圍是否與其他設定重疊（排除當前設定）
        overlapping = ScoreSettings.query.filter(
            ScoreSettings.id != setting_id,
            ScoreSettings.is_active == True,
            ((ScoreSettings.min_score <= max_score) & (ScoreSettings.max_score >= min_score))
        ).first()
        
        if overlapping and data.get('is_active', setting.is_active):
            return jsonify({'error': f'分數範圍與「{overlapping.level_name}」重疊'}), 400
        
        # 更新設定
        setting.level_name = data.get('level_name', setting.level_name)
        setting.min_score = min_score
        setting.max_score = max_score
        setting.description = data.get('description', setting.description)
        setting.is_active = data.get('is_active', setting.is_active)
        setting.order = data.get('order', setting.order)
        setting.updated_at = datetime.utcnow()
        
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '評分設定更新成功'
        })
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'更新評分設定失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/score-settings', methods=['POST'])
def create_score_setting():
    """創建新的評分設定"""
    try:
        data = request.json
        if not data:
            return jsonify({'error': '無效的請求數據'}), 400
        
        # 驗證必填欄位
        required_fields = ['level_name', 'min_score', 'max_score']
        for field in required_fields:
            if field not in data:
                return jsonify({'error': f'缺少必填欄位: {field}'}), 400
        
        min_score = data['min_score']
        max_score = data['max_score']
        
        if min_score >= max_score:
            return jsonify({'error': '最低分數必須小於最高分數'}), 400
        
        # 檢查分數範圍是否與其他啟用的設定重疊
        if data.get('is_active', True):
            overlapping = ScoreSettings.query.filter(
                ScoreSettings.is_active == True,
                ((ScoreSettings.min_score <= max_score) & (ScoreSettings.max_score >= min_score))
            ).first()
            
            if overlapping:
                return jsonify({'error': f'分數範圍與「{overlapping.level_name}」重疊'}), 400
        
        # 獲取下一個排序號
        max_order = db.session.query(func.max(ScoreSettings.order)).scalar() or 0
        
        # 創建新設定
        new_setting = ScoreSettings(
            level_name=data['level_name'],
            min_score=min_score,
            max_score=max_score,
            description=data.get('description', ''),
            is_active=data.get('is_active', True),
            order=data.get('order', max_order + 1)
        )
        
        db.session.add(new_setting)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '評分設定創建成功',
            'setting_id': new_setting.id
        })
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'創建評分設定失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/score-settings/<int:setting_id>', methods=['DELETE'])
def delete_score_setting(setting_id):
    """刪除評分設定"""
    try:
        setting = ScoreSettings.query.get(setting_id)
        if not setting:
            return jsonify({'error': '評分設定不存在'}), 404
        
        db.session.delete(setting)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '評分設定刪除成功'
        })
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'刪除評分設定失敗: {str(e)}'}), 500

def get_user_level_by_score(score):
    """根據分數獲取用戶等級"""
    try:
        setting = ScoreSettings.query.filter(
            ScoreSettings.is_active == True,
            ScoreSettings.min_score <= score,
            ScoreSettings.max_score >= score
        ).first()
        
        if setting:
            return setting.level_name
        else:
            # 如果沒有匹配的設定，返回默認值
            return '未分類'
            
    except Exception as e:
        print(f"獲取用戶等級失敗: {str(e)}")
        return '未分類'



# ==================== 推薦設定管理 API ====================

@quiz_bp.route('/api/admin/recommendation-settings', methods=['GET'])
def get_recommendation_settings():
    """獲取推薦設定列表"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    try:
        settings = RecommendationSettings.query.all()
        return jsonify([{
            'id': s.id,
            'setting_name': s.setting_name,
            'min_courses': s.min_courses,
            'max_courses': s.max_courses,
            'is_active': s.is_active,
            'description': s.description,
            'created_at': s.created_at.isoformat() if s.created_at else None,
            'updated_at': s.updated_at.isoformat() if s.updated_at else None
        } for s in settings])
    except Exception as e:
        return jsonify({'error': f'獲取推薦設定失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/recommendation-settings', methods=['POST'])
def create_recommendation_setting():
    """創建新的推薦設定"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    try:
        data = request.json
        
        # 驗證必要字段
        if not data.get('setting_name'):
            return jsonify({'error': '設定名稱不能為空'}), 400
        
        min_courses = data.get('min_courses', 3)
        max_courses = data.get('max_courses', 8)
        
        # 驗證數量範圍
        if min_courses < 0 or max_courses > 100:
            return jsonify({'error': '課程數量必須在0-100範圍內'}), 400
        
        if min_courses > max_courses:
            return jsonify({'error': '最少課程數量不能大於最多課程數量'}), 400
        
        # 檢查設定名稱是否已存在
        existing = RecommendationSettings.query.filter_by(setting_name=data['setting_name']).first()
        if existing:
            return jsonify({'error': '設定名稱已存在'}), 400
        
        # 創建新設定
        new_setting = RecommendationSettings(
            setting_name=data['setting_name'],
            min_courses=min_courses,
            max_courses=max_courses,
            is_active=data.get('is_active', True),
            description=data.get('description', '')
        )
        
        db.session.add(new_setting)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '推薦設定創建成功',
            'setting': {
                'id': new_setting.id,
                'setting_name': new_setting.setting_name,
                'min_courses': new_setting.min_courses,
                'max_courses': new_setting.max_courses,
                'is_active': new_setting.is_active,
                'description': new_setting.description
            }
        })
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'創建推薦設定失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/recommendation-settings/<int:setting_id>', methods=['PUT'])
def update_recommendation_setting(setting_id):
    """更新推薦設定"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    try:
        setting = RecommendationSettings.query.get_or_404(setting_id)
        data = request.json
        
        # 更新字段
        if 'setting_name' in data:
            # 檢查新名稱是否與其他設定衝突
            existing = RecommendationSettings.query.filter(
                RecommendationSettings.setting_name == data['setting_name'],
                RecommendationSettings.id != setting_id
            ).first()
            if existing:
                return jsonify({'error': '設定名稱已存在'}), 400
            setting.setting_name = data['setting_name']
        
        if 'min_courses' in data:
            min_courses = data['min_courses']
            if min_courses < 0 or min_courses > 100:
                return jsonify({'error': '最少課程數量必須在0-100範圍內'}), 400
            setting.min_courses = min_courses
        
        if 'max_courses' in data:
            max_courses = data['max_courses']
            if max_courses < 0 or max_courses > 100:
                return jsonify({'error': '最多課程數量必須在0-100範圍內'}), 400
            setting.max_courses = max_courses
        
        # 驗證範圍
        if setting.min_courses > setting.max_courses:
            return jsonify({'error': '最少課程數量不能大於最多課程數量'}), 400
        
        if 'is_active' in data:
            setting.is_active = data['is_active']
        
        if 'description' in data:
            setting.description = data['description']
        
        setting.updated_at = datetime.utcnow()
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '推薦設定更新成功',
            'setting': {
                'id': setting.id,
                'setting_name': setting.setting_name,
                'min_courses': setting.min_courses,
                'max_courses': setting.max_courses,
                'is_active': setting.is_active,
                'description': setting.description
            }
        })
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'更新推薦設定失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/recommendation-settings/<int:setting_id>', methods=['DELETE'])
def delete_recommendation_setting(setting_id):
    """刪除推薦設定"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    try:
        setting = RecommendationSettings.query.get_or_404(setting_id)
        
        # 不允許刪除默認設定
        if setting.setting_name == 'default':
            return jsonify({'error': '不能刪除默認設定'}), 400
        
        db.session.delete(setting)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '推薦設定刪除成功'
        })
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'刪除推薦設定失敗: {str(e)}'}), 500

@quiz_bp.route('/api/admin/recommendation-settings/<int:setting_id>/activate', methods=['POST'])
def activate_recommendation_setting(setting_id):
    """啟用推薦設定（同時停用其他設定）"""
    if not session.get('admin_logged_in'):
        return jsonify({'error': '未登錄'}), 401
    
    try:
        # 停用所有設定
        RecommendationSettings.query.update({'is_active': False})
        
        # 啟用指定設定
        setting = RecommendationSettings.query.get_or_404(setting_id)
        setting.is_active = True
        setting.updated_at = datetime.utcnow()
        
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': f'推薦設定 "{setting.setting_name}" 已啟用',
            'active_setting': {
                'id': setting.id,
                'setting_name': setting.setting_name,
                'min_courses': setting.min_courses,
                'max_courses': setting.max_courses
            }
        })
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'啟用推薦設定失敗: {str(e)}'}), 500

def get_active_recommendation_setting():
    """獲取當前啟用的推薦設定"""
    try:
        active_setting = RecommendationSettings.query.filter_by(is_active=True).first()
        if active_setting:
            return {
                'min_courses': active_setting.min_courses,
                'max_courses': active_setting.max_courses,
                'setting_name': active_setting.setting_name
            }
        else:
            # 如果沒有啟用的設定，返回默認值
            return {
                'min_courses': 3,
                'max_courses': 8,
                'setting_name': 'default'
            }
    except Exception as e:
        print(f"獲取推薦設定失敗: {str(e)}")
        return {
            'min_courses': 3,
            'max_courses': 8,
            'setting_name': 'default'
        }

